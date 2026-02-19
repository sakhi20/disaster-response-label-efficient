"""
Generate Project_Progress_Feb19.pptx
Structure: Heilmeier's Catechism
Run: python scripts/make_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)   # dark navy background
ACCENT     = RGBColor(0x00, 0x8B, 0xD8)   # bright blue accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xD0, 0xD8, 0xE8)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
ORANGE     = RGBColor(0xF3, 0x9C, 0x12)
DARK_BLUE  = RGBColor(0x0A, 0x1A, 0x35)
SLIDE_W    = Inches(13.33)
SLIDE_H    = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=16, color=WHITE, bullet="▸ ", indent_items=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    indent_items = indent_items or set()

    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        prefix = "    ◦ " if i in indent_items else bullet
        run = p.add_run()
        run.text = prefix + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def slide_header(slide, title):
    """Full-sentence title stating the message of the slide."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ACCENT)
    add_text(slide, title,
             Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.0),
             font_size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(1.15), Inches(12.33), Inches(0.03), ACCENT)


# ── Build Presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# SLIDE 1: Title
s1 = prs.slides.add_slide(blank)
set_bg(s1, NAVY)
add_rect(s1, 0, 0, Inches(0.12), SLIDE_H, ACCENT)
add_text(s1, "Label-Efficient Disaster Response\nUsing Foundation Models",
         Inches(1.0), Inches(2.0), Inches(11.33), Inches(2.0),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "Sakhi Patel   ·   Vivek Vanera   ·   Mulya Patel",
         Inches(1.0), Inches(4.2), Inches(11.33), Inches(0.5),
         font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s1, "February 19, 2026",
         Inches(1.0), Inches(4.8), Inches(11.33), Inches(0.4),
         font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)


# SLIDE 2: What are you trying to do?
s2 = prs.slides.add_slide(blank)
set_bg(s2, NAVY)
slide_header(s2, "The Problem: Disasters Strike, But Labels Take Days")

# Placeholder for visuals (Hurricane, Wildfire, Flood)
y_img = Inches(1.5)
h_img = Inches(3.0)
w_img = Inches(3.8)
margin = Inches(0.4)

types = ["Hurricane Damage", "Wildfire Damage", "Flood Damage"]
for i, t in enumerate(types):
    x = Inches(0.5 + i * (w_img + margin))
    add_rect(s2, x, y_img, w_img, h_img, DARK_BLUE)
    add_text(s2, t, x, y_img - Inches(0.3), w_img, Inches(0.3), font_size=14, bold=True, color=ORANGE)
    add_text(s2, "Before", x, y_img + h_img/2, w_img/2, Inches(0.3), font_size=12, align=PP_ALIGN.CENTER)
    add_text(s2, "After", x + w_img/2, y_img + h_img/2, w_img/2, Inches(0.3), font_size=12, align=PP_ALIGN.CENTER)
    add_text(s2, "VISUAL PLACEHOLDER", x, y_img + Inches(1.2), w_img, Inches(0.4), font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Simple statement
msg = (
    "When hurricanes hit, we have satellite photos in hours.\n"
    "But mapping damage requires human experts and takes days.\n"
    "Emergency response can't wait."
)
add_text(s2, msg, Inches(0.5), Inches(5.0), Inches(12.33), Inches(2.0),
         font_size=20, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


# SLIDE 3: How is it done today?
s3 = prs.slides.add_slide(blank)
set_bg(s3, NAVY)
slide_header(s3, "Current Approach: Deep Learning Needs Thousands of Labels")
add_bullet_box(s3, [
    "Traditional deep learning requires 100,000+ labeled examples to be effective.",
    "Limitation: Experts cannot label images fast enough during an active disaster.",
    "Limitation: Models often need to be retrained for every new disaster type or location.",
], Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0), font_size=22)


# SLIDE 4: What's new in our approach?
s4 = prs.slides.add_slide(blank)
set_bg(s4, NAVY)
slide_header(s4, "Our Solution: Foundation Models Learn Earth's 'Grammar'")
add_bullet_box(s4, [
    "Phase 1: Foundation models (Prithvi) learn patterns from unlabeled satellite data.",
    "Phase 2: Use these pre-learned features to detect damage with just 10% of labels.",
    "Innovation: Knowledge transfers across disaster types and global geographies.",
], Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0), font_size=22)


# SLIDE 5: Technical Approach
s5 = prs.slides.add_slide(blank)
set_bg(s5, NAVY)
slide_header(s5, "Two-Phase Pipeline")
# Diagram placeholder
add_rect(s5, Inches(1.0), Inches(1.8), Inches(11.33), Inches(4.0), DARK_BLUE)
add_text(s5, "Prithvi (Satellite Pre-training)  →  Vision Transformer (Damage Classifier)",
         Inches(1.0), Inches(3.0), Inches(11.33), Inches(1.0), font_size=24, bold=True, align=PP_ALIGN.CENTER)
add_text(s5, "Dataset: xBD (850,000 buildings, 19 natural disasters)",
         Inches(1.0), Inches(4.5), Inches(11.33), Inches(1.0), font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# SLIDE 6: Experimental Setup
s6 = prs.slides.add_slide(blank)
set_bg(s6, NAVY)
slide_header(s6, "Current Progress: Week 3")
add_bullet_box(s6, [
    "✓ GitHub repository established for collaborative research.",
    "✓ xBD dataset (51GB) downloaded and currently being extracted.",
    "✓ Baseline training script ready (ResNet-18) for benchmarking.",
    "Next: Fine-tuning geospatial foundation model (Prithvi).",
], Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0), font_size=22)


# SLIDE 7: Who cares?
s7 = prs.slides.add_slide(blank)
set_bg(s7, NAVY)
slide_header(s7, "Impact: Save Lives Through Faster Response")
add_bullet_box(s7, [
    "Faster damage assessment → Faster deployment of food, medical aid, and rescue teams.",
    "Global Generalization: One model works globally without targeted retraining.",
    "Reduces labeling burden by 90%: Strong results even when expert labels are scarce.",
], Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0), font_size=22)


# SLIDE 8: Future Directions
s8 = prs.slides.add_slide(blank)
set_bg(s8, NAVY)
slide_header(s8, "Future Directions: SAR Integration")
add_bullet_box(s8, [
    "Extend to Synthetic Aperture Radar (SAR) for all-weather disaster monitoring.",
    "SAR works through clouds and darkness — critical during active storms.",
    "Collaborative research with Prof. Mikhail Gilman (NCSU) on SAR geometry.",
], Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0), font_size=22)


# SLIDE 9: Conclusions
s9 = prs.slides.add_slide(blank)
set_bg(s9, NAVY)
slide_header(s9, "Target: 90%+ Accuracy with 10% Labels")
add_bullet_box(s9, [
    "We are building AI that 'understands' satellite data before a disaster even happens.",
    "Goal: High-precision damage detection with minimal human intervention.",
    "Targeting: CVPR EarthVision 2026 for research dissemination.",
], Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0), font_size=22)


prs.save("Project_Progress_Feb19.pptx")
print("✓ Saved: Project_Progress_Feb19.pptx")
